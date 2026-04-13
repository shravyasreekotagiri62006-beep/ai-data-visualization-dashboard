from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Apply a standard Master template if available, fallback is plain
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # TITLE SLIDE
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI Data Visualization Dashboard"
    subtitle.text = "Predictive Analytics, Automated Cleaning, and Dynamic Graphing\n\nBTech Final Year Project"

    # SLIDE 1: Introduction
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "1. Introduction"
    tf = body_shape.text_frame
    tf.text = "Overview of the Dashboard Project"
    p = tf.add_paragraph()
    p.text = "A full-stack web intelligent platform built specifically to ingest, secure, and analyze raw datasets (CSV/JSON) automatically."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Bridges the gap between data engineering and business intelligence."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Generates human-readable 'Mock AI' context including anomalies, summaries, and mathematical predictive trajectories natively."
    p.level = 1

    # SLIDE 2: Existing System
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "2. Existing System"
    tf = body_shape.text_frame
    tf.text = "Challenges in current Data Analysis workflows:"
    p = tf.add_paragraph()
    p.text = "Reliance on expensive enterprise tools (Tableau, PowerBI) requiring high technical overhead."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Users must manually clean missing or malformed Data rows in Excel before visualizations can map properly."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "No out-of-the-box system interprets the statistics contextually directly within the platform (e.g., standard deviations or outliers usually require manual coding)."
    p.level = 1

    # SLIDE 3: Proposed System
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "3. Proposed System"
    tf = body_shape.text_frame
    tf.text = "The AI Data Viz Dashboard revolutionizes this process:"
    p = tf.add_paragraph()
    p.text = "Automated Data Cleaning: Missing numerical nodes are natively inputted via Pandas means without breaking the app."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Machine Learning Integration: Scikit-learn K-Means unsupervised modeling instantly groups data components into statistical clusters."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Drag & Drop UI: Highly custom GridStack DOM coupled with Chart.js allowing modular graphic building horizontally."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Predictive Capabilities using NumPy polyfit trajectories."
    p.level = 1

    # SLIDE 4: System Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "4. System Architecture"
    tf = body_shape.text_frame
    tf.text = "The Tech Stack Foundation:"
    p = tf.add_paragraph()
    p.text = "Data Layer: Local lightweight SQLite relational database managing User Authentication (werkzeug hashed keys) and historically uploaded dataset paths."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Processor Engine: Python Flask combined with Pandas and NumPy mathematically isolating standard deviations and producing AI summaries."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Client Interface: Pure standard CSS Glassmorphism logic interacting with advanced Fetch APIs connected to interactive dynamically plotted Chart components."
    p.level = 1

    # SLIDE 5: Explanation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "5. Explanation & Workflow"
    tf = body_shape.text_frame
    tf.text = "How a user interacts with the system:"
    p = tf.add_paragraph()
    p.text = "1. Secure Login & Upload: User authenticates and uploads a CSV. Pandas traps exceptions to stop system breakdown on malformed data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Analysis Engine Triggers: NumPy predicts future trajectory bounds while KMeans calculates centroids."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Render Dashboard: AI insights display seamlessly without popups, and user visually constructs vibrant Bar/Line matrices."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Export: User downloads the ML-cleaned dataset or generates a high-quality PDF DOM Snapshot."
    p.level = 1

    prs.save('Project_Presentation.pptx')
    print("PPT created successfully")

create_presentation()
